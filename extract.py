# -*- coding: utf-8 -*-
"""Extract clean, RAG-ready plain text from Confluence attachments.
Allowlist: pdf, docx, pptx, txt, md, csv. Images/zip/scripts are skipped."""
import io
import re


def _pdf(b):
    import pdfplumber
    out = []
    with pdfplumber.open(io.BytesIO(b)) as pdf:
        for pg in pdf.pages:
            t = pg.extract_text() or ""
            if t.strip():
                out.append(t)
    return "\n".join(out)


def _docx(b):
    import docx
    d = docx.Document(io.BytesIO(b))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for tbl in d.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _pptx(b):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(b))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [sh.text for sh in slide.shapes
                 if sh.has_text_frame and sh.text and sh.text.strip()]
        if texts:
            parts.append(f"[Slide {i}] " + "  ".join(texts))
    return "\n".join(parts)


def _plain(b):
    return b.decode("utf-8", "ignore")


ALLOW = {"pdf": _pdf, "docx": _docx, "pptx": _pptx,
         "txt": _plain, "md": _plain, "csv": _plain}

# icon-font glyphs (Unicode private-use area) + control chars — noise for embeddings.
_PUA = re.compile("[%s-%s]" % (chr(0xE000), chr(0xF8FF)))
_CTRL = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")       # keep \n and \t


def _clean(t):
    t = t.replace("\r", "\n")
    t = _PUA.sub("", t)
    t = _CTRL.sub("", t)
    t = re.sub(r"[ \t]+", " ", t)
    t = re.sub(r"\n[ \t]+", "\n", t)
    t = re.sub(r"\n{3,}", "\n\n", t)
    return t.strip()


def extract_text(filename, data):
    """Return (clean_text, ext, ok). ok=False for unsupported types or failures."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    fn = ALLOW.get(ext)
    if not fn:
        return "", ext, False
    try:
        text = fn(data)
    except Exception:
        return "", ext, False
    return _clean(text), ext, True
